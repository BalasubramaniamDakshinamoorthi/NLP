import os
import datetime
import time
import re
import sys
import platform
import docx
import requests
import spacy
import uuid
import json
from io import BytesIO
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.pdfpage import PDFPage
import pptx

#Parse a list of paragraphs
def parse_and_insert_paragraph_list(paragraph_list):

    sentence_list = [] # A list of strings
    paragraph_number = 0
    sentence_number = 0
    word_number = 0
    paragraph_token_list = []
    sentence_token_list = []
    word_token_list = []
    sentence_paragraph_mapping = {}

    document_id = str(uuid.uuid4())

    # Create paragraph_token_list

    for paragraph in paragraph_list:

        if re.search('[A-Za-z0-9]', paragraph):

            print('paragraph = ' + paragraph)
            paragraph_token = {}
            paragraph_token['document_id'] = document_id
            paragraph_token['filename'] = filename
            paragraph_token['datatype'] = 'ascii'
            paragraph_token['paragraph_number'] = str(paragraph_number)
            paragraph_token['paragraph_text'] = paragraph
            paragraph_number += 1
            paragraph_token_list.append(paragraph_token)


    # Create sentence token list with references to paragraph numbers

    paragraph_number = 0
    sentence_number = 0

    for paragraph in paragraph_list:

        if re.search('[A-Za-z0-9]', paragraph):

            print('paragraph = ' + paragraph)

            parsed_paragraph = nlp(paragraph)

            for sentence in parsed_paragraph.sents:

                print('sentence = ' + str(sentence) + '\n')
                sentence_list.append(str(sentence))
                sentence_token = {}
                sentence_token['document_id'] = document_id
                sentence_token['filename'] = filename
                sentence_token['datatype'] = 'ascii'
                sentence_token['paragraph_number'] = str(paragraph_number)
                sentence_token['sentence_number'] = str(sentence_number)
                sentence_token['sentence_text'] = str(sentence)
                sentence_paragraph_mapping[sentence_number] = paragraph_number
                sentence_number += 1
                sentence_token_list.append(sentence_token)

    # print('sentence_list = ' + str(sentence_list))
    # print('sentence_token_list = ' + str(sentence_token_list))

    # with open('sentence_token_list.json', 'w') as outfile:
        # json.dump(sentence_token_list, outfile, indent=4)

            paragraph_number += 1

    sentence_number = 0
    word_number = 0

    for sentence in sentence_list:

        print('sentence = ' + sentence)

        word_list = nlp(sentence)

        for word in word_list:

            # print('word = ' + word.text)
            # print('lemma = ' + word.lemma_)
            # print('POS = ' + word.pos_)

            word_token = {}

            if word.pos_ == 'NOUN' \
                    or word.pos_ == 'PROPN' \
                    or word.pos_ == 'VERB' \
                    or word.pos_ == 'ADJ' \
                    or  word.pos_ == 'ADV':

                word_token['document_id'] = document_id
                word_token['filename'] = filename
                word_token['datatype'] = 'ascii'
                word_token['paragraph_number'] = str(sentence_paragraph_mapping[sentence_number])
                word_token['sentence_number'] = str(sentence_number)
                word_token['word_number'] = str(word_number)
                word_token['word_text'] = word.text
                word_token['word_lemma'] = word.lemma_
                if word.pos_ == 'PROPN':
                    word_token['word_pos'] = 'noun' # hack for WordNet
                else:
                    word_token['word_pos'] = word.pos_.lower()

                word_token_list.append(word_token)

                word_number += 1 

        sentence_number += 1

    # with open('word_token_list.json', 'w') as outfile:
        # json.dump(word_token_list, outfile, indent=4)
    
    u = 'http://localhost:7201/graphdb/insertdata'
    h = {'Content-Type':'application/json'}

    # How many elements each chunk should have 
    n = 1000

    # using list comprehension 
    chunked_paragraph_token_list = [paragraph_token_list[i:i + n] \
            for i in range(0, len(paragraph_token_list), n)] 

    # using list comprehension 
    chunked_sentence_token_list = [sentence_token_list[i:i + n] \
            for i in range(0, len(sentence_token_list), n)] 

    # using list comprehension 
    chunked_word_token_list = [word_token_list[i:i + n] \
            for i in range(0, len(word_token_list), n)] 

    if paragraph_token_list != []:
        for chunk in chunked_paragraph_token_list:
            print('Inserting chunk = ' + str(chunk))
            r = requests.post(u, headers=h, json=chunk)

    if sentence_token_list != []:
        for chunk in chunked_sentence_token_list:
            print('Inserting chunk = ' + str(chunk))
            r = requests.post(u, headers=h, json=chunk)

    if word_token_list != []:
        for chunk in chunked_word_token_list:
            print('Inserting chunk = ' + str(chunk))
            r = requests.post(u, headers=h, json=chunk)
 
#Load text file
def getTextContent(filename):

    line_list = []

    with open(filename, 'r') as f:

        line_list = f.read().splitlines()

    print('Line list = ' + str(line_list))
    return line_list

#Extract text from DOCX
def getDocxContent(filename):

    doc = docx.Document(filename)

    paragraphList = [] 

    for para in doc.paragraphs:
        paragraphText = para.text
        paragraphText = paragraphText.encode('ascii', errors='ignore').decode()
        print("Paragraph text= " + paragraphText)
        if para.text != '':
            paragraphList.append(paragraphText)
    print('ParagraphList = ' + str(paragraphList))
    return paragraphList

#Extract text from PDF
def getPdfContent(filename):

    """
    Extract text from PDF file, and return
    the string contained inside
 
    :param path (str) path to the .pdf file
    :return: text (str) string extracted
    """
 
    rsrcmgr = PDFResourceManager()
    retstr = BytesIO()
    device = TextConverter(rsrcmgr, retstr)
    with open(filename, "rb") as fp:  # open in 'rb' mode to read PDF bytes
        interpreter = PDFPageInterpreter(rsrcmgr, device)
        for page in PDFPage.get_pages(fp, check_extractable=True):
            interpreter.process_page(page)
        device.close()
        text = retstr.getvalue()
        retstr.close()
    page = text.decode('utf-8')
    page = page.encode('ascii', errors='ignore').decode()

    pageList = [page]

    print('pageList = ' + str(pageList))

    return pageList
	

# Extract text from PowerPoint presentation.

def getPptxContent(filename):

    pres = pptx.Presentation(filename)

    fullText = ""
    paragraphList = [] 
    shapeList = []
    shapeText = ''

    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                shapeText = shape.text
                shapeText = \
                    shapeText.encode('ascii', errors='ignore').decode()
                print("Shape text= " + shapeText)
                if shapeText != '':
                    shapeList.append(shapeText)
    print('Shape List = ' + str(shapeList))
    return shapeList

file_dict = {}

cmdargs = sys.argv

cmdargs.pop(0) # Get rid of the script name, leaving only the actual arguments.

print(str(cmdargs))

for arg in cmdargs:

    if arg.endswith('.txt'):

       bare_filename = arg[:-4]
       print(bare_filename)
       file_dict[bare_filename] = arg 

    if arg.endswith('.docx'):

       bare_filename = arg[:-5]
       print(bare_filename)
       file_dict[bare_filename] = arg 

    if arg.endswith('.pptx'):

       bare_filename = arg[:-5]
       print(bare_filename)
       file_dict[bare_filename] = arg 

    if arg.endswith('.pdf'):

       bare_filename = arg[:-4]
       print(bare_filename)
       file_dict[bare_filename] = arg 

    print(arg)

print('file_dict = ' + str(file_dict))

token_list = []

token = {}

nlp = spacy.load('en_core_web_sm')

for bare_filename in file_dict:

    filename = file_dict[bare_filename]
    print('filename = ' + filename)

    paragraph_list = []

    if filename.endswith('.txt'):
        paragraph_list  = getTextContent(filename)
        print('Have .txt')

    if filename.endswith('.docx'):
        paragraph_list  = getDocxContent(filename)
        print('Have .docx')

    if filename.endswith('.pptx'):
        paragraph_list = getPptxContent(filename)
        print('Have .pptx')

    if filename.endswith('.pdf'):
        paragraph_list = getPdfContent(filename)
        print('Have .pdf')

    parse_and_insert_paragraph_list(paragraph_list)